from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────
GREEN       = RGBColor(0x2E, 0x7D, 0x32)   # deep green
GREEN_LIGHT = RGBColor(0x43, 0xA0, 0x47)   # accent green
DARK        = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xF4, 0xF6, 0xF8)
MUTED       = RGBColor(0x78, 0x90, 0x9C)
DONE_BG     = RGBColor(0xE8, 0xF5, 0xE9)
DONE_FG     = RGBColor(0x1B, 0x5E, 0x20)
WIP_BG      = RGBColor(0xFF, 0xF8, 0xE1)
WIP_FG      = RGBColor(0xE6, 0x5C, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text_box(slide, text, l, t, w, h,
                 size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, items, l, t, w, h,
                   size=16, color=DARK, spacing=1.15,
                   bullet="▸ ", done_items=None, wip_items=None):
    """items: list of strings. done_items / wip_items: sets for coloring."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after  = Pt(2)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        if done_items and item in done_items:
            run.font.color.rgb = DONE_FG
            run.font.bold = True
        elif wip_items and item in wip_items:
            run.font.color.rgb = WIP_FG
            run.font.bold = True
        else:
            run.font.color.rgb = color
    return txb

def title_bar(slide, title, subtitle=None):
    """Green left accent bar + title text."""
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, GREEN)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.7),
                 size=30, bold=True, color=GREEN)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4),
                     size=14, color=MUTED)
    # thin divider
    div = slide.shapes.add_shape(1, Inches(0.4), Inches(1.3), Inches(12.5), Pt(2))
    div.line.fill.background()
    div.fill.solid()
    div.fill.fore_color.rgb = GREEN_LIGHT

def badge(slide, text, l, t, w, h, bg, fg):
    add_rect(slide, l, t, w, h, bg)
    add_text_box(slide, text, l, t, w, h,
                 size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# Full green background top half
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(4.2), GREEN)
add_rect(s, Inches(0), Inches(4.2), SLIDE_W, Inches(3.3), GRAY)

# Logo placeholder circle
add_rect(s, Inches(0.6), Inches(0.55), Inches(1.1), Inches(1.1), GREEN_LIGHT)
add_text_box(s, "🌿", Inches(0.6), Inches(0.55), Inches(1.1), Inches(1.1),
             size=32, align=PP_ALIGN.CENTER, color=WHITE)

add_text_box(s, "AgroGebeya",
             Inches(1.9), Inches(0.5), Inches(10), Inches(0.7),
             size=18, bold=False, color=RGBColor(0xC8,0xE6,0xC9))

add_text_box(s, "Agro-Retail Management System",
             Inches(1.9), Inches(1.1), Inches(10), Inches(0.9),
             size=36, bold=True, color=WHITE)

add_text_box(s, "Progress Presentation",
             Inches(1.9), Inches(2.0), Inches(10), Inches(0.6),
             size=22, color=RGBColor(0xC8,0xE6,0xC9))

add_text_box(s, "Connecting Farmers · Retailers · Government Authorities",
             Inches(1.9), Inches(2.7), Inches(10), Inches(0.5),
             size=14, color=RGBColor(0xA5,0xD6,0xA7))

# Bottom info
add_text_box(s, "Team: AgroGebeya Dev Team",
             Inches(0.6), Inches(4.6), Inches(5), Inches(0.4),
             size=13, color=MUTED)
add_text_box(s, "Date: April 2026",
             Inches(0.6), Inches(5.0), Inches(5), Inches(0.4),
             size=13, color=MUTED)
add_text_box(s, "Stack: FastAPI · PostgreSQL · Next.js · TypeScript",
             Inches(0.6), Inches(5.4), Inches(10), Inches(0.4),
             size=13, color=MUTED)

# ══════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Problem Statement", "What challenges does AgroGebeya solve?")

problems = [
    "Farmers lack direct access to retail markets — middlemen reduce profits",
    "No unified platform for product listing, pricing, and inventory management",
    "Manual order processes lead to errors, delays, and poor traceability",
    "Payment handling is fragmented — no integrated local payment support",
    "Transport coordination between farms and retailers is unstructured",
    "Government authorities have no real-time visibility into agricultural trade",
    "No centralized audit trail or compliance monitoring for agri-transactions",
]
add_bullet_box(s, problems,
               Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
               size=15, color=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 3 – Solution Overview
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Solution Overview", "A unified digital marketplace for Ethiopian agriculture")

# 3 columns
cols = [
    ("👨‍🌾  Farmers", ["List & manage products", "Set prices & inventory", "Receive orders & payments", "Request transport"]),
    ("🏪  Retailers", ["Browse & order products", "Track order status", "Manage payments via Chapa", "Communicate with farmers"]),
    ("🏛️  Authorities", ["Monitor all transactions", "Manage user verification", "Access audit logs", "Generate reports"]),
]
col_w = Inches(3.9)
for i, (title, pts) in enumerate(cols):
    lft = Inches(0.4 + i * 4.3)
    add_rect(s, lft, Inches(1.5), col_w, Inches(5.5), WHITE)
    add_text_box(s, title, lft + Inches(0.15), Inches(1.65), col_w - Inches(0.3), Inches(0.5),
                 size=15, bold=True, color=GREEN)
    add_bullet_box(s, pts, lft + Inches(0.15), Inches(2.2), col_w - Inches(0.3), Inches(4.5),
                   size=13, bullet="• ")

# ══════════════════════════════════════════════════════════
# SLIDE 4 – System Architecture
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "System Architecture", "Full-stack web application with RESTful API")

layers = [
    ("Frontend  (Next.js 14 · TypeScript · Tailwind CSS)",
     WHITE, GREEN, Inches(1.5)),
    ("REST API  (FastAPI · Python · JWT Auth · Pydantic)",
     WHITE, GREEN_LIGHT, Inches(2.6)),
    ("Business Logic  (Services · Validators · Encryption · Notifications)",
     WHITE, MUTED, Inches(3.7)),
    ("Database  (PostgreSQL · SQLAlchemy · Alembic Migrations)",
     WHITE, DARK, Inches(4.8)),
]
for text, bg, accent, top in layers:
    add_rect(s, Inches(1.0), top, Inches(11.0), Inches(0.75), bg)
    add_rect(s, Inches(1.0), top, Inches(0.12), Inches(0.75), accent)
    add_text_box(s, text, Inches(1.25), top, Inches(10.5), Inches(0.75),
                 size=14, bold=False, color=DARK)

# arrows between layers
for top in [Inches(2.35), Inches(3.45), Inches(4.55)]:
    add_text_box(s, "⬇", Inches(6.4), top, Inches(0.5), Inches(0.25),
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Side notes
add_text_box(s, "External Services",
             Inches(0.1), Inches(2.6), Inches(0.9), Inches(2),
             size=10, color=MUTED)
add_bullet_box(s, ["Chapa Payments", "File Storage", "Email / SMS"],
               Inches(9.8), Inches(2.5), Inches(3.2), Inches(2),
               size=11, color=MUTED, bullet="→ ")

# ══════════════════════════════════════════════════════════
# SLIDE 5 – Implementation Progress
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Implementation Progress", "Backend complete · Frontend in progress")

# Legend
badge(s, "✔  Completed", Inches(9.5), Inches(0.3), Inches(1.6), Inches(0.38), DONE_BG, DONE_FG)
badge(s, "⏳  In Progress", Inches(11.2), Inches(0.3), Inches(1.7), Inches(0.38), WIP_BG, WIP_FG)

# Backend column
add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.6), WHITE)
add_text_box(s, "⚙  Backend (FastAPI)",
             Inches(0.55), Inches(1.6), Inches(5.5), Inches(0.45),
             size=15, bold=True, color=GREEN)
be_done = [
    "Auth & JWT (register, login, refresh)",
    "Role-based access control (Farmer / Retailer / Admin)",
    "Product CRUD + image upload",
    "Order management & status tracking",
    "Chapa payment integration",
    "Transport request & approval",
    "National ID verification (front & back)",
    "Notifications & preferences",
    "Admin dashboard APIs",
    "Audit logging & backup",
    "Alembic DB migrations",
]
add_bullet_box(s, be_done,
               Inches(0.55), Inches(2.1), Inches(5.5), Inches(4.8),
               size=12, done_items=set(be_done), bullet="✔ ")

# Frontend column
add_rect(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.6), WHITE)
add_text_box(s, "🖥  Frontend (Next.js)",
             Inches(6.95), Inches(1.6), Inches(5.8), Inches(0.45),
             size=15, bold=True, color=GREEN)
fe_done = ["Auth pages (login, register)", "Dashboard layout", "Product listing & detail", "Profile & settings pages", "National ID verification flow", "Admin verification panel", "Notification preferences"]
fe_wip  = ["Order management UI", "Payment flow UI", "Transport request UI", "Real-time messaging", "Admin reports & analytics"]
add_bullet_box(s, fe_done,
               Inches(6.95), Inches(2.1), Inches(5.8), Inches(2.6),
               size=12, done_items=set(fe_done), bullet="✔ ")
add_bullet_box(s, fe_wip,
               Inches(6.95), Inches(4.5), Inches(5.8), Inches(2.4),
               size=12, wip_items=set(fe_wip), bullet="⏳ ")

# ══════════════════════════════════════════════════════════
# SLIDE 6 – Key Features
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Key Features", "Grouped by functional domain")

feature_groups = [
    ("👤  Identity & Access",   ["User registration & login", "JWT authentication", "Role-based permissions", "National ID verification"]),
    ("📦  Product & Inventory", ["Product listing & pricing", "Multi-image upload", "Inventory tracking", "Category management"]),
    ("🛒  Orders",              ["Place & modify orders", "Approval workflow", "Status tracking", "Order history"]),
    ("💳  Payments",            ["Chapa integration", "Transaction recording", "Payment verification", "Refund support"]),
    ("🚚  Transport",           ["Request submission", "Admin approval", "Driver assignment", "Delivery tracking"]),
    ("🔔  Admin & Monitoring",  ["Audit logging", "User management", "Backup & recovery", "Reports & analytics"]),
]

col_w = Inches(4.0)
row_h = Inches(2.6)
for idx, (title, pts) in enumerate(feature_groups):
    col = idx % 3
    row = idx // 3
    lft = Inches(0.35 + col * 4.3)
    top = Inches(1.5 + row * 2.8)
    add_rect(s, lft, top, col_w, row_h, WHITE)
    add_rect(s, lft, top, col_w, Inches(0.42), GREEN if row == 0 else GREEN_LIGHT)
    add_text_box(s, title, lft + Inches(0.1), top + Inches(0.04),
                 col_w - Inches(0.2), Inches(0.38),
                 size=12, bold=True, color=WHITE)
    add_bullet_box(s, pts, lft + Inches(0.1), top + Inches(0.5),
                   col_w - Inches(0.2), row_h - Inches(0.6),
                   size=11, bullet="• ")

# ══════════════════════════════════════════════════════════
# SLIDE 7 – System Workflow
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "System Workflow", "End-to-end transaction flow")

steps = [
    ("1", "Farmer\nRegisters", "Uploads National ID\nfor verification"),
    ("2", "Lists\nProducts", "Sets price, quantity\n& images"),
    ("3", "Retailer\nOrders", "Browses catalogue\n& places order"),
    ("4", "Payment\nvia Chapa", "Secure payment\nprocessed & recorded"),
    ("5", "Transport\nRequest", "Farmer/retailer\nrequests delivery"),
    ("6", "Admin\nMonitors", "Audit logs, reports\n& user management"),
]

step_w = Inches(1.8)
gap    = Inches(0.22)
start  = Inches(0.35)
top    = Inches(2.2)

for i, (num, title, sub) in enumerate(steps):
    lft = start + i * (step_w + gap)
    # circle
    add_rect(s, lft + Inches(0.55), top, Inches(0.7), Inches(0.7), GREEN)
    add_text_box(s, num, lft + Inches(0.55), top, Inches(0.7), Inches(0.7),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # box
    add_rect(s, lft, top + Inches(0.75), step_w, Inches(2.8), WHITE)
    add_text_box(s, title, lft, top + Inches(0.85), step_w, Inches(0.7),
                 size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text_box(s, sub, lft, top + Inches(1.55), step_w, Inches(1.8),
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # arrow (not after last)
    if i < len(steps) - 1:
        add_text_box(s, "→", lft + step_w + Inches(0.02), top + Inches(0.1),
                     gap + Inches(0.1), Inches(0.5),
                     size=18, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 8 – Challenges
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Challenges Faced", "Technical and design hurdles encountered")

challenges = [
    ("🔐  Security & Encryption",    "Encrypting National IDs at rest while keeping queries efficient"),
    ("💳  Payment Integration",       "Handling Chapa webhook callbacks and idempotent transaction recording"),
    ("📸  ID Verification UX",        "Building a 2-step flow (number entry + front/back photo upload) on mobile"),
    ("🔄  Async DB Operations",       "Managing SQLAlchemy async sessions with proper error handling and rollbacks"),
    ("🌐  CORS & Error Propagation",  "500 errors stripping CORS headers — fixed by moving middleware to top of stack"),
    ("📱  Mobile Responsiveness",     "Adapting complex settings and admin pages to small screen layouts"),
    ("🗂️  Migration Management",      "Coordinating Alembic migrations across multiple feature branches"),
]

for i, (title, desc) in enumerate(challenges):
    top = Inches(1.55 + i * 0.73)
    add_rect(s, Inches(0.4), top, Inches(12.4), Inches(0.62), WHITE)
    add_text_box(s, title, Inches(0.55), top + Inches(0.08),
                 Inches(3.2), Inches(0.45), size=12, bold=True, color=GREEN)
    add_text_box(s, desc, Inches(3.8), top + Inches(0.08),
                 Inches(8.8), Inches(0.45), size=12, color=DARK)

# ══════════════════════════════════════════════════════════
# SLIDE 9 – Next Steps
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GRAY)
title_bar(s, "Next Steps", "Roadmap to full completion")

short = [
    "Complete order management UI (placement, tracking, history)",
    "Build payment flow UI with Chapa redirect & confirmation",
    "Implement transport request & tracking screens",
    "Finish real-time messaging with WebSocket integration",
]
medium = [
    "Admin analytics dashboard with charts and export",
    "Mobile app (React Native) for farmers in the field",
    "SMS / push notification delivery integration",
    "End-to-end testing and performance optimisation",
]

add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.45), GREEN)
add_text_box(s, "🚀  Immediate (Next 2 Weeks)",
             Inches(0.55), Inches(1.54), Inches(5.6), Inches(0.38),
             size=13, bold=True, color=WHITE)
add_bullet_box(s, short, Inches(0.55), Inches(2.05), Inches(5.6), Inches(4.7),
               size=13, bullet="→ ")

add_rect(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.5), WHITE)
add_rect(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.45), GREEN_LIGHT)
add_text_box(s, "📅  Medium Term (1–2 Months)",
             Inches(7.05), Inches(1.54), Inches(5.7), Inches(0.38),
             size=13, bold=True, color=WHITE)
add_bullet_box(s, medium, Inches(7.05), Inches(2.05), Inches(5.7), Inches(4.7),
               size=13, bullet="→ ")

# ══════════════════════════════════════════════════════════
# SLIDE 10 – Conclusion
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, GREEN)
add_rect(s, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), DARK)

add_text_box(s, "🌿  AgroGebeya", Inches(1.5), Inches(0.5), Inches(10), Inches(0.7),
             size=18, color=RGBColor(0xC8,0xE6,0xC9), align=PP_ALIGN.CENTER)
add_text_box(s, "Bridging the Gap Between\nFarmers and Markets",
             Inches(1.0), Inches(1.2), Inches(11), Inches(1.4),
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

summary = [
    "✔  Fully implemented backend with 15+ API modules",
    "✔  Secure auth, payments, verification & transport",
    "⏳  Frontend ~60% complete — core flows in progress",
    "🎯  Production-ready architecture, scalable design",
]
add_bullet_box(s, summary, Inches(2.5), Inches(2.9), Inches(8.5), Inches(2.0),
               size=15, color=WHITE, bullet="")

add_text_box(s, "Thank You  ·  Questions Welcome",
             Inches(1.0), Inches(5.4), Inches(11), Inches(0.6),
             size=18, bold=True, color=RGBColor(0xA5,0xD6,0xA7), align=PP_ALIGN.CENTER)
add_text_box(s, "github.com/agrogebeya  ·  support@agrogebeya.com",
             Inches(1.0), Inches(6.1), Inches(11), Inches(0.45),
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────
out = "AgroGebeya_Progress_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
